from flask import Flask, jsonify, send_file
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_DATA_LABEL_POSITION
import plotly.express as px
import pandas as pd
import io


app = Flask(__name__)

@app.get('/generate')
def generate():
    try:
        prs = Presentation()
        title_layout_one = prs.slide_layouts[6]

        slide1 = prs.slides.add_slide(title_layout_one)
        slide2 = prs.slides.add_slide(title_layout_one)
        slide3 = prs.slides.add_slide(title_layout_one)
        slide4 = prs.slides.add_slide(title_layout_one)
        slide5 = prs.slides.add_slide(title_layout_one)
        slide6 = prs.slides.add_slide(title_layout_one)

        # line chart data
        line_chart_data = CategoryChartData()
        line_chart_data.categories = ['2014', '2015', '2016', '2017', '2018', '2019', '2020', '2021', '2022', '2023']
        line_chart_data.add_series('', (9000, 11000, 9500, 9500, 13000, 19000, 17000, 20000, 16000, 18000))
        # line chart size
        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)

        avt_chart = slide1.shapes.add_chart(
            XL_CHART_TYPE.LINE_STACKED, x, y, cx, cy, line_chart_data
        ).chart
        avt_chart.has_legend = False

        plot = avt_chart.plots[0]
        series = plot.series[0]
        line = series.format.line
        line.color.rgb = RGBColor(0, 100, 0)

        # Column Stacked
        stacked_bar_Chart_Data = CategoryChartData()
        stacked_bar_Chart_Data.categories = ['2014', '2015', '2016', '2017', '2018', '2019', '2020', '2021', '2022', '2023']
        stacked_bar_Chart_Data.add_series('WOL', (9000, 11000, 9500, 9500, 13000, 19000, 17000, 20000, 16000, 18000))
        stacked_bar_Chart_Data.add_series('EBSCO', (500, 600, 600, 600, 500, 500, 500, 500, 500, 500))

        x, y, cx, cy = Inches(2), Inches(2), Inches(7), Inches(4.5)

        stacked_bar_Chart = slide2.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_STACKED, x, y, cx, cy, stacked_bar_Chart_Data
        ).chart
        stacked_bar_Chart.has_legend = True
        stacked_bar_Chart.legend.include_in_layout = False
        stack_plot = stacked_bar_Chart.plots[0]
        seriesWOL = stack_plot.series[0]
        wolfill = seriesWOL.format.fill
        wolfill.solid()
        wolfill.fore_color.rgb = RGBColor(0, 100, 0)
        seriesEBSCO = stack_plot.series[1]
        ebscofill = seriesEBSCO.format.fill
        ebscofill.solid()
        ebscofill.fore_color.rgb = RGBColor(0, 0, 100)
        # column chart
        column_chart_data = CategoryChartData()
        column_chart_data.categories = ['2014', '2015', '2016', '2017', '2018', '2019', '2020', '2021', '2022', '2023']
        column_chart_data.add_series('', (9000, 11000, 9500, 9500, 13000, 19000, 17000, 20000, 16000, 18000))
        # line chart size
        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
        column_chart = slide3.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, column_chart_data
        ).chart
        column_chart.has_legend = False
        column_chart_plot = column_chart.plots[0]
        column_chart_series = column_chart_plot.series[0]
        columnCFill = column_chart_series.format.fill
        columnCFill.solid()
        columnCFill.fore_color.rgb = RGBColor(0, 100, 0)
        # pie chart
        pie_chart_data = ChartData()
        pie_chart_data.categories = ['West', 'East', 'North', 'South', 'Other']
        pie_chart_data.add_series('', (0.135, 0.324, 0.180, 0.235, 0.126))

        pie_chart = slide4.shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, pie_chart_data
        ).chart
        pie_chart.has_legend = True
        pie_chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        pie_chart.legend.include_in_layout = False

        pie_chart.plots[0].has_data_labels = True
        data_labels = pie_chart.plots[0].data_labels
        data_labels.number_format = '0%'
        data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

        # progress Chart with bar chart
        progress_chart_data = ChartData()
        progress_chart_data.categories = ['', ]
        progress_chart_data.add_series('Ad Impressions', (3100,))
        progress_chart = slide5.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, Inches(2), progress_chart_data
        ).chart
        progress_chart.overlap = -50

        # world map
        data = {
            'country': ['United States', 'Canada', 'Brazil', 'France', 'Germany', 'China', 'Australia'],
            'value': [100, 80, 70, 90, 85, 95, 60]
        }
        df = pd.DataFrame(data)

        custom_color_scale = [
            [0.0, 'rgb(220, 220, 255)'],
            [0.2, 'rgb(160, 160, 255)'],
            [0.4, 'rgb(100, 100, 255)'],
            [0.6, 'rgb(40, 40, 255)'],
            [0.8, 'rgb(0, 0, 220)'],
            [1.0, 'rgb(0, 0, 160)']
        ]
        # Create Choropleth map
        fig = px.choropleth(
            df,
            locations='country',
            locationmode='country names',
            color='value',
            color_continuous_scale=custom_color_scale,

        )
        fig.update_layout(
            coloraxis_showscale=False,
            margin=dict(l=0, r=0, t=0, b=0)  # Remove margins
        )
        img_bytes = fig.to_image(format="png", width=1920, height=1080)
        img_stream = io.BytesIO(img_bytes)
        img_stream.seek(0)
        choropleth_map = slide6.shapes.add_picture(img_stream, x, y, cx, cy)

        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)

        return send_file(pptx_io, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation', as_attachment=True, download_name='presentation.pptx')
    
    except Exception as e:
        return jsonify({"error": str(e)})

if __name__ == '__main__':
    app.run(debug=True)
